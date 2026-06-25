"""
Générateur de présentations PowerPoint LMS ORH.
Utilise python-pptx pour créer un benchmark RH mission au format LMS.

Design inspiré du template PHASE 1 Version 20260226.pptx :
  - Palette navy (#12294D) / bleu (#156082) / teal (#1B8D92)
  - Font Century Gothic
  - Sidebar gris gauche + header fin + callout So what

Fonction principale : generate_lms_ppt(analysis, mission_config, output_path) -> str
"""
import re
from datetime import datetime
from pathlib import Path
from typing import Dict, Any, List, Optional

try:
    from pptx import Presentation
    from pptx.util import Cm, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches
except ImportError as _e:
    raise ImportError(
        "python-pptx n'est pas installé. Exécutez : pip install python-pptx>=1.0.0"
    ) from _e

# ─────────────────────────────────────────────────────────────────────────────
#  CONSTANTES VISUELLES — PALETTE LMS "PHASE 1"
# ─────────────────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x12, 0x29, 0x4D)   # #12294D — bleu nuit principal
NAVY_DARK = RGBColor(0x0A, 0x18, 0x30)   # version plus sombre pour couvertures
BLUE      = RGBColor(0x15, 0x60, 0x82)   # #156082 — bleu accent (cercles)
TEAL      = RGBColor(0x1B, 0x8D, 0x92)   # #1B8D92 — teal secondaire
NAVY_PALE = RGBColor(0xCA, 0xDC, 0xFC)   # #CADCFC — bleu pâle (texte sur fond navy)
GRIS_BAR  = RGBColor(0xD8, 0xD8, 0xD8)   # #D8D8D8 — barre latérale grise
GRIS_FOND = RGBColor(0xF4, 0xF6, 0xF9)   # fond de contenu légèrement bleuté
GRIS_SEP  = RGBColor(0xD0, 0xD8, 0xE8)   # séparateur colonne
GRIS_DARK = RGBColor(0x2D, 0x2D, 0x2D)   # corps de texte sombre
BLANC     = RGBColor(0xFF, 0xFF, 0xFF)
TEAL_PALE = RGBColor(0xE5, 0xF4, 0xF5)   # fond callout So what
POLICE    = "Century Gothic"
POLICE_BODY = "Century Gothic"

# Dimensions layout (16:9)
HEADER_H  = Cm(1.4)    # bandeau titre mince (navy)
FOOTER_H  = Cm(0.9)    # bande confidentiel seulement (réduit — so what dans callout)
SO_WHAT_H = Cm(1.75)   # hauteur du callout so what (au-dessus footer)
BAR_V_W   = Cm(0.80)   # barre latérale grise gauche

SLIDE_W = Cm(33.87)
SLIDE_H = Cm(19.05)

# Assets
_HERE   = Path(__file__).parent
_ASSETS = _HERE / "assets"


def _find_asset(*candidates) -> Optional[Path]:
    for name in candidates:
        p = _ASSETS / name
        if p.exists():
            return p
    return None


LOGO_PATH   = _find_asset("logo_lms.png", "Logo LMS.png", "logo_lms.jpg", "Logo LMS.jpg")
MOSAIC_PATH = _find_asset("mosaic_lms.jpg", "mosaic_lms.png", "photo mosaïque.png",
                           "photo mosaique.png", "mosaic_lms.png")


# ─────────────────────────────────────────────────────────────────────────────
#  HELPERS DE BASE
# ─────────────────────────────────────────────────────────────────────────────

def _add_rect(slide, left, top, width, height, fill_color=None, line_color=None,
              rounded=False):
    """Ajoute un rectangle coloré."""
    if rounded:
        shape = slide.shapes.add_shape(5, left, top, width, height)
        shape.adjustments[0] = 0.05
    else:
        shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def _as_dict(val) -> dict:
    """Protège contre un LLM qui renvoie une string au lieu d'un dict."""
    return val if isinstance(val, dict) else {}


def _as_list(val) -> list:
    """Protège contre un LLM qui renvoie autre chose qu'une liste."""
    return val if isinstance(val, list) else []


def _clean_slide_text(text: str) -> str:
    """Supprime marqueurs internes ([confirmé], [N]) et nettoie les espaces."""
    if not isinstance(text, str):
        return ""
    text = re.sub(r'\[confirm[eé]\]|\[probable\]|\[[àa]\s*v[eé]rifier\]', '', text, flags=re.IGNORECASE)
    text = re.sub(r'(\[\d+\])+', '', text)
    text = re.sub(r'\s+([.,;:!?])', r'\1', text)
    text = re.sub(r'  +', ' ', text)
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()


def _short(text: str, max_chars: int = 200) -> str:
    """Tronque à la dernière phrase complète avant max_chars."""
    text = _clean_slide_text(text)
    if len(text) <= max_chars:
        return text
    cut = text[:max_chars]
    for sep in ('. ', '.\n', ' — ', ' : '):
        idx = cut.rfind(sep)
        if idx > max_chars // 2:
            return cut[:idx + 1]
    return cut.rstrip() + "…"


def _add_textbox(slide, left, top, width, height, text, font_size=12, bold=False,
                 color=None, align=PP_ALIGN.LEFT, wrap=True, italic=False,
                 font_name=None):
    """Ajoute une zone de texte simple."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name or POLICE
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox


def _add_numbered_circle(slide, left, top, number, size=Cm(0.75),
                          bg_color=None, text_color=BLANC, font_size=14):
    """Cercle numéroté style PHASE 1."""
    if bg_color is None:
        bg_color = BLUE if (number % 2 == 1) else TEAL
    circle = slide.shapes.add_shape(9, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = bg_color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(number)
    run.font.name = POLICE
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = text_color
    return circle


# ─────────────────────────────────────────────────────────────────────────────
#  HELPERS VISUELS RESTRUCTURÉS
# ─────────────────────────────────────────────────────────────────────────────

def _split_to_bullets(text: str, max_items: int = 5, max_chars_per: int = 90) -> List[str]:
    """
    Découpe un texte paragraphe en liste de bullets courts.
    Chaque phrase terminée (. ! ?) suivie d'une majuscule = 1 bullet.
    """
    text = _clean_slide_text(text)
    if not text:
        return []
    # Split after sentence-ending punctuation followed by uppercase
    parts = re.split(
        r'(?<=[.!?])\s+(?=[A-ZÀÂÄÈÉÊËÎÏÔÙÛÜŸŒÆ])|(?<=[.!?])\n|\n{1,2}',
        text
    )
    result = []
    for p in parts:
        p = p.strip().rstrip('.')
        if len(p) < 10:
            continue
        if len(p) > max_chars_per:
            cut = p[:max_chars_per].rsplit(' ', 1)[0]
            p = cut + '…'
        result.append(p)
        if len(result) >= max_items:
            break
    return result


def _add_bullet_list_box(slide, left, top, width, height, items: List[str],
                          header_color=None, item_color=GRIS_DARK,
                          header_size: float = 10, item_size: float = 10,
                          arrow_color=None):
    """
    Rendu d'une liste à puces structurée.
    Convention items :
      '§Titre'    → sous-titre bold (header_color)
      'texte'     → puce '→ texte' (arrow en teal, texte en item_color)
      '' / None   → ignoré
    """
    if header_color is None:
        header_color = BLUE
    if arrow_color is None:
        arrow_color = TEAL

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    first_para = True
    for raw in items:
        if not raw or not str(raw).strip():
            continue
        raw = _clean_slide_text(str(raw))
        if not raw:
            continue

        p = tf.paragraphs[0] if first_para else tf.add_paragraph()
        first_para = False

        if raw.startswith('§'):
            # Sous-titre de section
            p.space_before = Pt(5)
            run = p.add_run()
            run.text = raw[1:].strip()
            run.font.name = POLICE
            run.font.size = Pt(header_size)
            run.font.bold = True
            run.font.color.rgb = header_color
        else:
            # Puce avec flèche colorée
            p.space_before = Pt(1)
            run_arrow = p.add_run()
            run_arrow.text = "→ "
            run_arrow.font.name = POLICE
            run_arrow.font.size = Pt(item_size)
            run_arrow.font.bold = False
            run_arrow.font.color.rgb = arrow_color

            run_text = p.add_run()
            run_text.text = raw
            run_text.font.name = POLICE
            run_text.font.size = Pt(item_size)
            run_text.font.bold = False
            run_text.font.color.rgb = item_color

    return txBox


def _add_so_what_callout(slide, left, top, width, height, entreprise: str, text: str):
    """
    Callout teal clair pour le So what — intégré dans la zone de contenu.
    [accent bar teal] ▶ Pour ENTREPRISE : texte
    """
    # Fond teal très clair
    _add_rect(slide, left, top, width, height, fill_color=TEAL_PALE)
    # Barre d'accent gauche teal
    _add_rect(slide, left, top, Cm(0.22), height, fill_color=TEAL)

    sw_text = _short(text or "", max_chars=210)

    txBox = slide.shapes.add_textbox(
        left + Cm(0.38), top + Cm(0.08),
        width - Cm(0.55), height - Cm(0.1)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]

    # "▶ Pour ENTREPRISE : " en bold teal
    run1 = p.add_run()
    run1.text = f"▶ Pour {entreprise} : "
    run1.font.name = POLICE
    run1.font.size = Pt(9)
    run1.font.bold = True
    run1.font.color.rgb = TEAL

    # Texte So what en navy
    run2 = p.add_run()
    run2.text = sw_text
    run2.font.name = POLICE
    run2.font.size = Pt(9)
    run2.font.bold = False
    run2.font.color.rgb = NAVY


# ─────────────────────────────────────────────────────────────────────────────
#  CHROME COMMUN (header + sidebar + footer confidentiel)
# ─────────────────────────────────────────────────────────────────────────────

def _add_slide_chrome(slide, titre: str, nom_mission: str, entreprise_cible: str) -> float:
    """
    Pose header + sidebar + footer confidentiel.
    Retourne footer_top (float Emu) pour que les slides calculent leur zone de contenu.
    Le So what est désormais un callout dans la zone de contenu, pas dans le footer.
    """
    # Sidebar gris gauche (pleine hauteur)
    _add_rect(slide, Cm(0), Cm(0), BAR_V_W, SLIDE_H, fill_color=GRIS_BAR)

    # Header navy
    _add_rect(slide, BAR_V_W, Cm(0), SLIDE_W - BAR_V_W, HEADER_H, fill_color=NAVY)
    _add_textbox(
        slide,
        BAR_V_W + Cm(0.4), Cm(0.1),
        SLIDE_W - BAR_V_W - Cm(5.0), HEADER_H - Cm(0.1),
        titre, font_size=15, bold=True, color=BLANC,
    )

    # Logo LMS dans le header
    if LOGO_PATH:
        try:
            logo_w = Cm(3.8)
            logo_h = HEADER_H - Cm(0.1)
            slide.shapes.add_picture(
                str(LOGO_PATH),
                SLIDE_W - logo_w - Cm(0.2), Cm(0.05),
                logo_w, logo_h,
            )
        except Exception:
            pass

    # Footer navy mince (confidentiel seulement)
    footer_top = SLIDE_H - FOOTER_H
    _add_rect(slide, BAR_V_W, footer_top, SLIDE_W - BAR_V_W, FOOTER_H, fill_color=NAVY)
    _add_textbox(
        slide,
        SLIDE_W - Cm(11.5), footer_top + Cm(0.1),
        Cm(11.0), FOOTER_H - Cm(0.1),
        f"Confidentiel — {nom_mission}  ·  LMS ORH",
        font_size=7.5, color=NAVY_PALE, align=PP_ALIGN.RIGHT,
    )

    return footer_top


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE DE CONTENU RESTRUCTURÉ
# ─────────────────────────────────────────────────────────────────────────────

def _content_slide(prs, titre: str,
                   col_gauche_items: List[str], col_droite_items: List[str],
                   so_what_text: str, mission_config: dict,
                   col_gauche_titre: str = "Analyse",
                   col_droite_titre: str = "Benchmark"):
    """
    Slide de contenu structuré — bullets + callout so_what.

    Layout :
    ┌──┬─────────────────────────────────────────────────────────┐
    │▌ │ HEADER NAVY — titre                        [LMS logo] │
    ├──┼───────────────────────┬─────────────────────────────────┤
    │▌ │ col_gauche_titre      │ col_droite_titre               │
    │▌ │  → bullet 1           │  § Section                    │
    │▌ │  → bullet 2           │  → KPI 1 : valeur             │
    │▌ │  § Sous-section       │  → KPI 2 : valeur             │
    │▌ │  → bullet 3           │                               │
    ├──┴───────────────────────────────────────────────────────── │
    │▌ [TEAL CALLOUT] ▶ Pour Entreprise : So what              │
    ├──┴─────────────────────────────────────────────────────────┤
    │▌ [FOOTER] Confidentiel — Mission · LMS ORH               │
    └──────────────────────────────────────────────────────────── ┘
    """
    nom_mission      = mission_config.get("nom_mission", "Mission")
    entreprise_cible = mission_config.get("entreprise_cible", "Entreprise")

    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    sp_tree = slide.shapes._spTree
    if len(sp_tree) > 2:
        sp_tree.remove(sp_tree[2])

    footer_top = _add_slide_chrome(slide, titre, nom_mission, entreprise_cible)

    content_top = HEADER_H + Cm(0.2)
    # Zone blanche totale (header → footer)
    _add_rect(slide, BAR_V_W, content_top,
              SLIDE_W - BAR_V_W, footer_top - content_top, fill_color=BLANC)

    # Callout So what juste au-dessus du footer
    sw_top = footer_top - SO_WHAT_H
    _add_so_what_callout(
        slide,
        BAR_V_W, sw_top,
        SLIDE_W - BAR_V_W, SO_WHAT_H,
        entreprise_cible, so_what_text
    )

    # Zone bullets (entre header et callout so_what)
    bullet_zone_h = sw_top - content_top - Cm(0.15)

    # Colonnes
    inner_left = BAR_V_W + Cm(0.55)
    col_l_w    = (SLIDE_W - BAR_V_W - Cm(1.3)) * 0.56
    sep_x      = inner_left + col_l_w + Cm(0.15)
    col_r_x    = sep_x + Cm(0.28)
    col_r_w    = SLIDE_W - col_r_x - Cm(0.45)

    # Séparateur vertical
    _add_rect(slide, sep_x, content_top + Cm(0.2), Cm(0.04),
              bullet_zone_h - Cm(0.3), fill_color=GRIS_SEP)

    # ── Colonne gauche ────────────────────────────────────────────────────────
    _add_textbox(
        slide, inner_left, content_top + Cm(0.1), col_l_w, Cm(0.5),
        col_gauche_titre, font_size=10, bold=True, color=BLUE,
    )
    _add_bullet_list_box(
        slide, inner_left, content_top + Cm(0.65),
        col_l_w, bullet_zone_h - Cm(0.7),
        col_gauche_items or [],
        header_color=BLUE, item_color=GRIS_DARK, arrow_color=TEAL,
        header_size=10, item_size=10,
    )

    # ── Colonne droite ────────────────────────────────────────────────────────
    _add_textbox(
        slide, col_r_x, content_top + Cm(0.1), col_r_w, Cm(0.5),
        col_droite_titre, font_size=10, bold=True, color=TEAL,
    )
    _add_bullet_list_box(
        slide, col_r_x, content_top + Cm(0.65),
        col_r_w, bullet_zone_h - Cm(0.7),
        col_droite_items or [],
        header_color=TEAL, item_color=GRIS_DARK, arrow_color=BLUE,
        header_size=10, item_size=10,
    )

    return slide


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDES SPÉCIFIQUES
# ─────────────────────────────────────────────────────────────────────────────

def _slide_cover(prs, analysis, mission_config):
    """Slide 1 — Couverture style PHASE 1."""
    nom_mission      = mission_config.get("nom_mission", "Mission RH")
    entreprise_cible = mission_config.get("entreprise_cible", "Entreprise")
    secteur          = mission_config.get("secteur", "")
    geographie       = mission_config.get("geographie", "")
    concurrent       = mission_config.get("concurrent_reference", "") or ""

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Panneau gauche ~33%
    panneau_w = SLIDE_W * 0.33
    _add_rect(slide, Cm(0), Cm(0), panneau_w, SLIDE_H, fill_color=NAVY_DARK)
    if MOSAIC_PATH:
        try:
            slide.shapes.add_picture(str(MOSAIC_PATH), Cm(0), Cm(0), panneau_w, SLIDE_H)
        except Exception:
            pass
    _add_rect(slide, Cm(0), SLIDE_H * 0.65, panneau_w, SLIDE_H * 0.35, fill_color=NAVY_DARK)
    _add_textbox(
        slide, Cm(0.3), SLIDE_H - Cm(2.0), panneau_w - Cm(0.4), Cm(1.6),
        "LMS ORH", font_size=16, bold=True, color=BLANC, align=PP_ALIGN.CENTER,
    )

    # Zone droite fond blanc
    right_x = panneau_w + Cm(0.5)
    right_w = SLIDE_W - right_x - Cm(0.6)
    _add_rect(slide, panneau_w, Cm(0), SLIDE_W - panneau_w, SLIDE_H, fill_color=BLANC)

    _bench_type_label = mission_config.get("type", "RH").upper()
    _add_textbox(
        slide, right_x, Cm(1.8), right_w, Cm(0.7),
        f"BENCHMARK {_bench_type_label} — MISSION CONSULTANT",
        font_size=10, bold=True, color=BLUE,
    )
    _add_rect(slide, right_x, Cm(2.65), right_w * 0.6, Cm(0.06), fill_color=NAVY)
    _add_textbox(
        slide, right_x, Cm(3.0), right_w, Cm(2.5),
        nom_mission, font_size=26, bold=True, color=NAVY,
    )
    _add_textbox(
        slide, right_x, Cm(5.8), right_w, Cm(1.3),
        entreprise_cible, font_size=20, bold=True, color=BLUE,
    )
    sect_geo = f"{secteur}  |  {geographie}" if geographie else secteur
    _add_textbox(slide, right_x, Cm(7.3), right_w, Cm(0.8), sect_geo, font_size=12, color=GRIS_DARK)

    if concurrent:
        _add_textbox(
            slide, right_x, Cm(8.3), right_w, Cm(0.7),
            f"Référence : {concurrent}", font_size=11, color=RGBColor(0x60, 0x80, 0xA0),
        )
        date_top = Cm(9.2)
    else:
        date_top = Cm(8.3)

    date_str = datetime.now().strftime("%d %B %Y")
    _add_textbox(slide, right_x, date_top, right_w, Cm(0.7), date_str, font_size=11, color=GRIS_DARK)
    _add_textbox(
        slide, right_x, SLIDE_H - Cm(2.0), right_w, Cm(0.6),
        "Document confidentiel — LMS ORH",
        font_size=9, color=RGBColor(0x90, 0xA0, 0xB0),
    )
    if LOGO_PATH:
        try:
            logo_w = Cm(5.0)
            logo_h = Cm(2.0)
            slide.shapes.add_picture(
                str(LOGO_PATH),
                SLIDE_W - logo_w - Cm(0.5), SLIDE_H - logo_h - Cm(0.3),
                logo_w, logo_h,
            )
        except Exception:
            pass
    return slide


def _slide_contexte(prs, analysis, mission_config):
    """Slide 2 — Contexte & Angle RH."""
    ctx = _as_dict(analysis.get("contexte_mission"))
    texte = ctx.get("texte", "") or ctx.get("angle_organisationnel", "")
    angle = ctx.get("angle_rh", "") or ctx.get("angle_organisationnel", "")
    angle_strategique = mission_config.get("angle_strategique_rh", "")
    geographie = mission_config.get("geographie", "")
    concurrent = mission_config.get("concurrent_reference", "") or ""

    left_items = _split_to_bullets(texte, max_items=5)

    right_items: List[str] = []
    if angle:
        right_items += ["§Angle central"] + _split_to_bullets(angle, max_items=2)
    if geographie:
        right_items += ["§Périmètre géographique", geographie]
    if concurrent:
        right_items += ["§Référence comparative", concurrent]
    if angle_strategique:
        right_items += ["§Question stratégique"] + _split_to_bullets(angle_strategique, max_items=2)

    so_what = angle_strategique or angle or "Voir angle RH central."
    return _content_slide(
        prs,
        titre="Contexte & Angle RH",
        col_gauche_items=left_items,
        col_droite_items=right_items,
        so_what_text=so_what,
        mission_config=mission_config,
        col_gauche_titre="Contexte mission",
        col_droite_titre="Cadrage stratégique",
    )


def _slide_business_model(prs, analysis, mission_config, titre=None):
    """Slide 3 — Business Model — Lecture RH."""
    bm = _as_dict(analysis.get("business_model_rh"))
    analyse    = bm.get("analyse", "")
    emergentes = _as_list(bm.get("competences_emergentes"))
    obsoletes  = _as_list(bm.get("competences_obsoletes"))
    kpis       = _as_list(bm.get("kpis_sectoriels"))
    so_what    = bm.get("so_what", "")

    left_items = _split_to_bullets(analyse, max_items=3)
    if emergentes:
        left_items += ["§Compétences émergentes"] + [str(c) for c in emergentes[:4]]

    right_items: List[str] = []
    if obsoletes:
        right_items += ["§Compétences en déclin"] + [str(c) for c in obsoletes[:3]]
    if kpis:
        right_items += ["§KPIs RH clés"] + [str(k) for k in kpis[:4]]

    return _content_slide(
        prs,
        titre=titre or "Business Model — Lecture RH",
        col_gauche_items=left_items,
        col_droite_items=right_items,
        so_what_text=so_what,
        mission_config=mission_config,
        col_gauche_titre="Analyse RH",
        col_droite_titre="Compétences & KPIs",
    )


def _slide_organisation(prs, analysis, mission_config, titre=None):
    """Slide 4 — Organisation & Dimensionnement."""
    org = _as_dict(analysis.get("organisation_dimensionnement"))
    nouveaux_roles = _as_list(org.get("nouveaux_roles"))
    kpis           = _as_list(org.get("kpis_sectoriels"))

    left_items = _split_to_bullets(org.get("analyse", ""), max_items=3)
    if nouveaux_roles:
        left_items += ["§Nouveaux rôles"] + [str(r) for r in nouveaux_roles[:3]]

    right_items: List[str] = []
    if kpis:
        right_items += ["§KPIs sectoriels"] + [str(k) for k in kpis[:3]]
    tendances = org.get("tendances_effectifs", "")
    if tendances:
        right_items += ["§Tendances effectifs"] + _split_to_bullets(tendances, max_items=2)
    externalisation = org.get("externalisation", "")
    if externalisation:
        right_items += ["§Externalisation"] + _split_to_bullets(externalisation, max_items=2)

    return _content_slide(
        prs,
        titre=titre or "Organisation & Dimensionnement",
        col_gauche_items=left_items,
        col_droite_items=right_items,
        so_what_text=org.get("so_what", ""),
        mission_config=mission_config,
        col_gauche_titre="Analyse",
        col_droite_titre="KPIs & Tendances",
    )


def _slide_gouvernance(prs, analysis, mission_config, titre=None):
    """Slide 5 — Gouvernance RH & Management."""
    gov  = _as_dict(analysis.get("gouvernance_rh"))
    kpis = _as_list(gov.get("kpis_sectoriels"))

    left_items = _split_to_bullets(gov.get("analyse", ""), max_items=4)

    right_items: List[str] = []
    if gov.get("instances_rh"):
        right_items += ["§Instances RH"] + _split_to_bullets(gov["instances_rh"], max_items=2)
    if gov.get("politiques_sociales"):
        right_items += ["§Politiques sociales"] + _split_to_bullets(gov["politiques_sociales"], max_items=2)
    if gov.get("conformite"):
        right_items += ["§Conformité"] + _split_to_bullets(gov["conformite"], max_items=1)
    if kpis:
        right_items += ["§KPIs gouvernance"] + [str(k) for k in kpis[:3]]

    return _content_slide(
        prs,
        titre=titre or "Gouvernance RH & Management",
        col_gauche_items=left_items,
        col_droite_items=right_items,
        so_what_text=gov.get("so_what", ""),
        mission_config=mission_config,
        col_gauche_titre="Analyse",
        col_droite_titre="Gouvernance & Conformité",
    )


def _slide_signaux_innovation(prs, analysis, mission_config, titre=None):
    """Slide 6 — Signaux faibles & Innovations RH."""
    inn     = _as_dict(analysis.get("innovation_manageriale"))
    signaux = _as_list(analysis.get("signaux_faibles"))
    kpis    = _as_list(inn.get("kpis_sectoriels"))

    left_items: List[str] = []
    for s in signaux[:4]:
        s = _as_dict(s)
        signal  = s.get("signal", "")
        horizon = s.get("horizon", "")
        if signal:
            left_items.append(f"{signal}" + (f" ({horizon})" if horizon else ""))
    if not left_items:
        left_items = _split_to_bullets(inn.get("analyse", ""), max_items=4)

    pratiques = _as_list(inn.get("pratiques_differenciantes"))
    outils    = _as_list(inn.get("outils_rh"))

    right_items: List[str] = []
    if pratiques:
        right_items += ["§Pratiques différenciantes"] + [str(p) for p in pratiques[:3]]
    if outils:
        right_items += ["§Outils RH"] + [str(o) for o in outils[:3]]
    if kpis:
        right_items += ["§KPIs innovation"] + [str(k) for k in kpis[:2]]

    return _content_slide(
        prs,
        titre=titre or "Signaux faibles & Innovations RH",
        col_gauche_items=left_items,
        col_droite_items=right_items,
        so_what_text=inn.get("so_what", ""),
        mission_config=mission_config,
        col_gauche_titre="Signaux faibles",
        col_droite_titre="Innovations & KPIs",
    )


# ── Mode ORGANISATIONNEL ──────────────────────────────────────────────────────

def _slide_modeles_csp(prs, analysis, mission_config, titre=None):
    """Slide 3 Org — Modèles CSP comparables."""
    csp        = _as_dict(analysis.get("modeles_csp"))
    structures = _as_list(csp.get("structures_types"))
    kpis       = _as_list(csp.get("kpis_sectoriels"))

    left_items = _split_to_bullets(csp.get("analyse", ""), max_items=3)
    if structures:
        left_items += ["§Structures observées"] + [str(s) for s in structures[:4]]

    right_items: List[str] = []
    if kpis:
        right_items += ["§KPIs CSP"] + [str(k) for k in kpis[:3]]
    if csp.get("gouvernance_observee"):
        right_items += ["§Gouvernance observée"] + _split_to_bullets(csp["gouvernance_observee"], max_items=2)
    if csp.get("perimetre_fonctionnel"):
        right_items += ["§Périmètre fonctionnel"] + _split_to_bullets(csp["perimetre_fonctionnel"], max_items=2)

    return _content_slide(
        prs,
        titre=titre or "Modèles CSP — Benchmark comparatif",
        col_gauche_items=left_items,
        col_droite_items=right_items,
        so_what_text=csp.get("so_what", ""),
        mission_config=mission_config,
        col_gauche_titre="Analyse comparative",
        col_droite_titre="KPIs & Gouvernance",
    )


def _slide_processus_douaniers(prs, analysis, mission_config, titre=None):
    """Slide 4 Org — Processus douaniers & import/export."""
    pd       = _as_dict(analysis.get("processus_douaniers"))
    pratiques = _as_list(pd.get("bonnes_pratiques"))
    outils    = _as_list(pd.get("outils_systemes"))
    risques   = _as_list(pd.get("risques_frequents"))
    kpis      = _as_list(pd.get("kpis_sectoriels"))

    left_items: List[str] = []
    if pratiques:
        left_items += [str(p) for p in pratiques[:5]]
    else:
        left_items = _split_to_bullets(pd.get("analyse", ""), max_items=5)

    right_items: List[str] = []
    if outils:
        right_items += ["§Outils & systèmes"] + [str(o) for o in outils[:3]]
    if risques:
        right_items += ["§Risques fréquents"] + [str(r) for r in risques[:3]]
    if kpis:
        right_items += ["§KPIs douaniers"] + [str(k) for k in kpis[:3]]

    return _content_slide(
        prs,
        titre=titre or "Processus douaniers — Best practices",
        col_gauche_items=left_items,
        col_droite_items=right_items,
        so_what_text=pd.get("so_what", ""),
        mission_config=mission_config,
        col_gauche_titre="Bonnes pratiques",
        col_droite_titre="Outils · Risques · KPIs",
    )


def _slide_interface_filiale_siege(prs, analysis, mission_config, titre=None):
    """Slide 5 Org — Interface filiale/siège."""
    iface = _as_dict(analysis.get("interface_filiale_siege"))
    kpis  = _as_list(iface.get("kpis_sectoriels"))

    left_items = _split_to_bullets(iface.get("analyse", ""), max_items=3)
    if iface.get("modeles_delegation"):
        left_items += ["§Modèles de délégation"] + _split_to_bullets(iface["modeles_delegation"], max_items=2)

    right_items: List[str] = []
    if iface.get("protocoles_validation"):
        right_items += ["§Protocoles validation"] + _split_to_bullets(iface["protocoles_validation"], max_items=2)
    if iface.get("reporting_type"):
        right_items += ["§Reporting"] + _split_to_bullets(iface["reporting_type"], max_items=2)
    if kpis:
        right_items += ["§KPIs interface"] + [str(k) for k in kpis[:3]]

    return _content_slide(
        prs,
        titre=titre or "Interface Filiale / Siège",
        col_gauche_items=left_items,
        col_droite_items=right_items,
        so_what_text=iface.get("so_what", ""),
        mission_config=mission_config,
        col_gauche_titre="Analyse",
        col_droite_titre="Protocoles & Reporting",
    )


def _slide_formalisation_audit(prs, analysis, mission_config, titre=None):
    """Slide 6 Org — Formalisation & Audit-readiness + grille maturité."""
    fau          = _as_dict(analysis.get("formalisation_audit_readiness"))
    signaux      = _as_list(analysis.get("signaux_faibles"))
    referentiels = _as_list(fau.get("referentiels_utilises"))
    criteres     = _as_list(fau.get("criteres_audit_groupe"))
    grille       = _as_list(fau.get("grille_maturite"))
    kpis         = _as_list(fau.get("kpis_sectoriels"))

    left_items: List[str] = []
    if referentiels:
        left_items += ["§Référentiels utilisés"] + [str(r) for r in referentiels[:3]]
    if criteres:
        left_items += ["§Critères audit groupe"] + [str(c) for c in criteres[:4]]
    if not left_items:
        left_items = _split_to_bullets(fau.get("analyse", ""), max_items=5)

    right_items: List[str] = []
    if grille:
        right_items += ["§Grille de maturité (world-class)"]
        for g in grille[:4]:
            g = _as_dict(g)
            axe = g.get("axe", "")
            pratiques = g.get("pratiques_world_class", "")
            if axe:
                if isinstance(pratiques, str) and pratiques:
                    pratiques_short = pratiques[:70] + "…" if len(pratiques) > 70 else pratiques
                    right_items.append(f"{axe} : {pratiques_short}")
                else:
                    right_items.append(axe)
    elif fau.get("niveaux_maturite"):
        right_items += ["§Niveaux de maturité"] + _split_to_bullets(fau["niveaux_maturite"], max_items=3)

    if kpis:
        right_items += ["§KPIs audit"] + [str(k) for k in kpis[:2]]

    # Signaux faibles org
    siglist: List[str] = []
    for s in signaux[:2]:
        s = _as_dict(s)
        signal = s.get("signal", "")
        if signal:
            siglist.append(str(signal))
    if siglist:
        right_items += ["§Signaux faibles"] + siglist

    return _content_slide(
        prs,
        titre=titre or "Formalisation & Audit-readiness",
        col_gauche_items=left_items,
        col_droite_items=right_items,
        so_what_text=fau.get("so_what", ""),
        mission_config=mission_config,
        col_gauche_titre="Référentiels & Critères",
        col_droite_titre="Maturité & Signaux",
    )


def _slide_recommandations(prs, analysis, mission_config):
    """
    Slide Recommandations — 3 blocs numérotés style PHASE 1.
    Cercles numérotés blue/teal + titre + justification + KPI + horizon.
    """
    nom_mission      = mission_config.get("nom_mission", "Mission")
    entreprise_cible = mission_config.get("entreprise_cible", "Entreprise")
    recs = _as_list(analysis.get("recommandations_mission"))

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    footer_top = _add_slide_chrome(slide, "Recommandations Mission", nom_mission, entreprise_cible)

    content_top = HEADER_H + Cm(0.3)
    content_h   = footer_top - content_top - Cm(0.2)
    _add_rect(slide, BAR_V_W, content_top, SLIDE_W - BAR_V_W, content_h, fill_color=BLANC)

    nb_recs  = min(len(recs), 3)
    if nb_recs == 0:
        return slide
    inner_w  = SLIDE_W - BAR_V_W - Cm(1.0)
    gap      = Cm(0.45)
    bloc_w   = (inner_w - gap * (nb_recs - 1)) / nb_recs
    bloc_top = content_top + Cm(0.25)
    bloc_h   = content_h - Cm(0.35)

    priorite_colors = {"Haute": NAVY, "Moyenne": BLUE, "Faible": TEAL}

    for i, rec in enumerate(recs[:3]):
        rec  = _as_dict(rec)
        left = BAR_V_W + Cm(0.45) + i * (bloc_w + gap)

        _add_rect(slide, left, bloc_top, bloc_w, bloc_h, fill_color=GRIS_FOND)

        # Cercle numéroté
        circle_size = Cm(0.9)
        circle_left = left + (bloc_w - circle_size) / 2
        _add_numbered_circle(
            slide, circle_left, bloc_top + Cm(0.2),
            number=i + 1, size=circle_size,
            bg_color=BLUE if i % 2 == 0 else TEAL,
        )

        # Titre de la reco
        action = _short(rec.get("action", f"Recommandation {i+1}"), max_chars=160)
        _add_textbox(
            slide, left + Cm(0.3), bloc_top + Cm(1.3),
            bloc_w - Cm(0.5), Cm(1.8),
            action, font_size=11, bold=True, color=NAVY,
        )

        # Justification — bullets si possible
        justif_raw = rec.get("justification", "")
        justif_bullets = _split_to_bullets(justif_raw, max_items=4, max_chars_per=80)
        if justif_bullets:
            _add_bullet_list_box(
                slide, left + Cm(0.3), bloc_top + Cm(3.25),
                bloc_w - Cm(0.5), Cm(5.5),
                justif_bullets,
                header_color=NAVY, item_color=GRIS_DARK, arrow_color=TEAL,
                item_size=9,
            )
        else:
            _add_textbox(
                slide, left + Cm(0.3), bloc_top + Cm(3.25),
                bloc_w - Cm(0.5), Cm(5.5),
                _short(justif_raw, max_chars=320),
                font_size=9, color=GRIS_DARK,
            )

        # Priorité
        priorite = _clean_slide_text(rec.get("priorite", "Moyenne"))
        p_color  = priorite_colors.get(priorite, NAVY)
        _add_textbox(
            slide, left + Cm(0.3), bloc_h + bloc_top - Cm(2.7),
            bloc_w - Cm(0.5), Cm(0.6),
            f"● Priorité : {priorite}", font_size=9, bold=True, color=p_color,
        )

        # KPI
        kpi = _short(rec.get("kpi", ""), max_chars=150)
        _add_textbox(
            slide, left + Cm(0.3), bloc_h + bloc_top - Cm(2.0),
            bloc_w - Cm(0.5), Cm(1.2),
            f"KPI : {kpi}", font_size=9, color=TEAL,
        )

        # Horizon
        horizon = _clean_slide_text(rec.get("horizon", ""))
        _add_textbox(
            slide, left + Cm(0.3), bloc_h + bloc_top - Cm(0.75),
            bloc_w - Cm(0.5), Cm(0.65),
            f"Horizon : {horizon}", font_size=9, color=GRIS_DARK,
        )

    return slide


def _slide_optionnelle(prs, slide_data, mission_config):
    """Slide optionnelle — format contenu standard."""
    titre       = slide_data.get("titre", "Analyse thématique")
    observation = slide_data.get("observation", "")
    benchmark   = slide_data.get("benchmark_sectoriel", "")
    so_what     = slide_data.get("so_what", "")

    left_items  = _split_to_bullets(observation, max_items=5)
    right_items = _split_to_bullets(benchmark, max_items=5)

    return _content_slide(
        prs,
        titre=titre,
        col_gauche_items=left_items,
        col_droite_items=right_items,
        so_what_text=so_what,
        mission_config=mission_config,
    )


# ─────────────────────────────────────────────────────────────────────────────
#  FONCTION PRINCIPALE
# ─────────────────────────────────────────────────────────────────────────────

def generate_lms_ppt(analysis: Dict[str, Any], mission_config: Dict, output_path: str) -> str:
    """
    Génère un fichier PowerPoint LMS ORH pour un benchmark mission.

    Args:
        analysis      : dict retourné par analyze_mission()
        mission_config: dict de configuration de la mission
        output_path   : chemin absolu du fichier .pptx à créer

    Returns:
        Le chemin du fichier généré (= output_path).
    """
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    _is_org = mission_config.get("type", "RH").upper() == "ORGANISATIONNEL"
    _axes   = mission_config.get("axes_noms") or {}

    if _is_org:
        _t1 = _axes.get("axe1") or "Modèles CSP — Benchmark comparatif"
        _t2 = _axes.get("axe2") or "Processus douaniers — Best practices"
        _t3 = _axes.get("axe3") or "Interface Filiale / Siège"
        _t4 = _axes.get("axe4") or "Formalisation & Audit-readiness"
    else:
        _t1 = _axes.get("axe1") or "Business Model — Lecture RH"
        _t2 = _axes.get("axe2") or "Organisation & Dimensionnement"
        _t3 = _axes.get("axe3") or "Gouvernance RH & Management"
        _t4 = _axes.get("axe4") or "Signaux faibles & Innovations RH"

    _slide_cover(prs, analysis, mission_config)
    _slide_contexte(prs, analysis, mission_config)
    if _is_org:
        _slide_modeles_csp(prs, analysis, mission_config, titre=_t1)
        _slide_processus_douaniers(prs, analysis, mission_config, titre=_t2)
        _slide_interface_filiale_siege(prs, analysis, mission_config, titre=_t3)
        _slide_formalisation_audit(prs, analysis, mission_config, titre=_t4)
    else:
        _slide_business_model(prs, analysis, mission_config, titre=_t1)
        _slide_organisation(prs, analysis, mission_config, titre=_t2)
        _slide_gouvernance(prs, analysis, mission_config, titre=_t3)
        _slide_signaux_innovation(prs, analysis, mission_config, titre=_t4)
    _slide_recommandations(prs, analysis, mission_config)

    slides_opt = analysis.get("slides_optionnelles", [])
    if isinstance(slides_opt, list):
        for slide_data in slides_opt:
            if isinstance(slide_data, dict):
                _slide_optionnelle(prs, slide_data, mission_config)

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return str(out)
